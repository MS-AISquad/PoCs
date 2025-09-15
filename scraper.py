import pypdf
from docx import Document
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.text.text import _Paragraph
from pptx.table import _Cell
import requests
from io import BytesIO


class Reader:
    def __init__(self, input_name: str):
        self.input_name = input_name
        self.doc = None
        self.paragraphs = None
        self.join_char = ' '
        if input_name.endswith('.pdf'):
            self.text = self.read_pdf(input_name)
            self.extension = 'pdf'
        elif input_name.endswith('.txt'):
            self.text = [self.read_txt(input_name)]
            self.extension = 'txt'
        elif input_name.endswith('.docx'):
            self.doc, self.paragraphs = self.read_docx(input_name)
            self.text = [p.text for p in self.paragraphs]
            self.join_char = '\n\n' + '#'*20 + '\n\n'
            self.extension = 'docx'
        elif input_name.endswith('.pptx'):
            self.doc, self.paragraphs = self.read_pptx(input_name)
            self.text = [p.text for p in self.paragraphs]
            self.join_char = '\n\n' + '#'*20 + '\n\n'
            self.extension = 'pptx'
        elif input_name.startswith('http'):
            self.text = [self.read_webpage(input_name)]
        else:
            raise ValueError(f'Path was not recognized: \n{input_name}')
        

    def read_pdf(self, file_path: str) -> list[str]:

        if file_path.startswith('http'):
            
            req = requests.get(file_path)
            assert req.status_code == 200, f'Request for pdf file was unsuccessful:\n{file_path}'
            reader = pypdf.PdfReader(BytesIO(req.content))

        else:

            # creating a pdf reader object
            reader = pypdf.PdfReader(file_path)
        
        text = [page.extract_text() for page in reader.pages]
        return text


    def read_txt(self, file_path):

        text = open(file_path, 'r').read()
        return text


    def read_webpage(self, url):

        req = requests.get(url)
        assert req.status_code == 200, f'Request for webpage file was unsuccessful:\n{url}'

        return req.text


    def read_docx(self, file_path: str) -> list[Paragraph]:

        doc = Document(file_path)

        paragraph_lst = []
        paragraph_lst.extend(doc.paragraphs)
        paragraph_lst.extend(paragraph 
                            for section in doc.sections 
                            for paragraph in section.header.paragraphs)
        paragraph_lst.extend(paragraph 
                            for section in doc.sections 
                            for paragraph in section.footer.paragraphs)
        paragraph_lst.extend(paragraph 
                            for table in doc.tables 
                            for row in table.table.rows 
                            for cell in row.cells 
                            for paragraph in cell.paragraphs)
        
        return doc, paragraph_lst
    

    def read_pptx(self, file_path: str) -> list[_Paragraph | _Cell]:

        prs = Presentation(file_path)

        paragraph_lst = []
        paragraph_lst.extend(paragraph
                            for slide in prs.slides
                            for shape in slide.shapes
                            if shape.has_text_frame
                            for paragraph in shape.text_frame.paragraphs)
        paragraph_lst.extend(cell
                            for slide in prs.slides
                            for shape in slide.shapes
                            if shape.has_table
                            for row in shape.table.rows
                            for cell in row.cells)
        paragraph_lst.extend(slide.notes_slide.notes_text_frame
                            for slide in prs.slides
                            if slide.has_notes_slide)
        
        return prs, paragraph_lst
    



# read_pdf(r'documents\Notice of Repurchasing shares\inputs\Notice of Repurchasing shares ENG.pdf')
# read_pdf('https://www.managementsolutions.com/sites/default/files/minisite/static/72b0015f-39c9-4a52-ba63-872c115bfbd0/llm/pdf/rise-of-llm.pdf')

# read_txt(r'documents\Notice of Repurchasing shares\inputs\Notice of Repurchasing shares ENG.txt')

# read_webpage('https://en.wikipedia.org/wiki/Large_language_model')


def test_translate(text):
    import re
    translation = []
    for word in text.split(' '):
        if word.lower().startswith('y'):
            i = re.search('[aeiouy]', word.lower()[1:])
            extra = 1
        else:
            i = re.search('[aeiouy]', word.lower())
            extra = 0
        if i is None:
            translation.append(word)
            continue
        else:
            i = i.start() + extra
        word = word[i:] + word[:i] + 'ay' if i > 0 else word + 'yay'
        word = re.sub(r'(\w+)([.,!?\'"])(\w+)', r'\1\3\2', word) 
        translation.append(word)
    return ' '.join(translation)


# below is an example of how translation can be done for docx and re-inserted with the same structure.
# note that run-level format/style is lost because we won't want to translate below sentence-level,
# so only paragraph-level format/style can be maintained
"""
read_doc = Reader(r'documents\Test document.docx')
doc, paragraphs = read_doc.doc, read_doc.text
for p in paragraphs:
    text = p.text
    # perform translation, below is a placeholder function:
    translated = test_translate(text)
    p.text = translated

doc.save(r'documents\test_output.docx')
"""

"""
read_prs = Reader(r'documents\Test Powerpoint.pptx')
doc, paragraphs, texts = read_prs.doc, read_prs.paragraphs, read_prs.text
for text, p in zip(texts, paragraphs):
    # perform translation, below is a placeholder function:
    translated = test_translate(text)
    p.text = translated

doc.save(r'documents\test_output.pptx')
"""
